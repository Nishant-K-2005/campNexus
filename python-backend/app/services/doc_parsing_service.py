import pymupdf
import io
from pptx import Presentation
from docx import Document
import pandas as pd

async def parseFiles(docFile, mime):
    mime = mime.split('/')[-1]
    if(mime=="pdf"):
        return await getTextFromPdf(docFile)
    mime = mime.split('.')[-1]
    if(mime=="presentation"):
        return await getTextFromPpt(docFile)
    if(mime=="document"):
        return await getTextFromDoc(docFile)
    if(mime=="csv" or mime=="sheet"):
        return await getTextFromSheet(docFile,mime)

async def getTextFromPdf(file):
    text = ""
    file.file.seek(0)
    pdf_bytes = await file.read()
    pdf = pymupdf.open(stream=pdf_bytes,filetype="pdf")
    for page in pdf:
        text += page.get_text()   
    pdf.close()
    return text

async def getTextFromPpt(file):
    text = ""
    file.file.seek(0)
    ppt_content = await file.read()
    ppt_stream = io.BytesIO(ppt_content)
    prs = Presentation(ppt_stream)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape,"text"):
                text += shape.text + "\n"

    return text

async def getTextFromDoc(file):
    text = ""
    file.file.seek(0)
    doc_content = await file.read()
    doc_stream = io.BytesIO(doc_content)
    doc = Document(doc_stream)
    for para in doc.paragraphs:
        if para.text:
            text+=para.text+"\n"
    return text

async def getTextFromSheet(file, mime):
    text=""
    file.file.seek(0)
    table_content = await file.read()
    table_stream = io.BytesIO(table_content)
    if(mime=="sheet"):
        df = pd.read_excel(table_stream)
    if(mime=="csv"):
        df = pd.read_csv(table_stream)
    text = df.to_string(index=False)
    return text+"\n"